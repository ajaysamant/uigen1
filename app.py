import os
from urllib.parse import urlparse

import truststore

truststore.inject_into_ssl()

import requests
import streamlit as st
from anthropic import Anthropic
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from pptx import Presentation
from pypdf import PdfReader
from docx import Document

load_dotenv()

MAX_CHARS = 150_000
MODEL = "claude-opus-4-8"

st.set_page_config(page_title="Document & URL Summarizer", page_icon="📄")
st.title("Document & URL Summarizer")


@st.cache_resource
def get_client() -> Anthropic:
    return Anthropic()


def extract_pdf(file) -> str:
    reader = PdfReader(file)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_docx(file) -> str:
    doc = Document(file)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_pptx(file) -> str:
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def extract_url(url: str) -> str:
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError("URL must start with http:// or https://")

    resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
    resp.raise_for_status()

    soup = BeautifulSoup(resp.text, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def summarize(client: Anthropic, content: str) -> str:
    response = client.messages.create(
        model=MODEL,
        max_tokens=4096,
        messages=[
            {
                "role": "user",
                "content": (
                    "Summarize the following content clearly and concisely, "
                    "covering the key points:\n\n" + content
                ),
            }
        ],
    )
    return next((b.text for b in response.content if b.type == "text"), "")


if not os.environ.get("ANTHROPIC_API_KEY"):
    st.error("ANTHROPIC_API_KEY is not set. Add it to your .env file and restart the app.")
    st.stop()

client = get_client()

source = st.radio("Choose input source", ["Upload document", "Paste URL"])

content = None
label = None

if source == "Upload document":
    uploaded = st.file_uploader(
        "Upload a PDF, Word, or PowerPoint file", type=["pdf", "docx", "pptx"]
    )
    if uploaded is not None:
        label = uploaded.name
        with st.spinner("Extracting text..."):
            try:
                name = uploaded.name.lower()
                if name.endswith(".pdf"):
                    content = extract_pdf(uploaded)
                elif name.endswith(".docx"):
                    content = extract_docx(uploaded)
                elif name.endswith(".pptx"):
                    content = extract_pptx(uploaded)
            except Exception as e:
                st.error(f"Could not read file: {e}")
else:
    url = st.text_input("Enter a public URL")
    if url and st.button("Fetch & Summarize"):
        try:
            with st.spinner("Fetching page..."):
                content = extract_url(url)
            label = url
        except Exception as e:
            st.error(f"Could not fetch URL: {e}")

if content is not None:
    if not content.strip():
        st.warning("No extractable text was found in this source.")
    else:
        if len(content) > MAX_CHARS:
            st.warning(
                f"Content is long ({len(content):,} characters) — only the first "
                f"{MAX_CHARS:,} characters will be summarized."
            )
            content = content[:MAX_CHARS]

        should_summarize = (
            source == "Paste URL" or st.button(f"Summarize {label}")
        )
        if should_summarize:
            with st.spinner("Summarizing with Claude..."):
                try:
                    summary = summarize(client, content)
                except Exception as e:
                    st.error(f"Summarization failed: {e}")
                else:
                    st.subheader("Summary")
                    st.write(summary)
